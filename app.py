"""
Enhanced PowerPoint Generator Backend - FIXED TEMPLATE PRESERVATION
The key fix is in create_enhanced_presentation function
"""

from flask import Flask, request, send_file, jsonify, render_template_string
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import io, json, os, requests, mimetypes, traceback
from werkzeug.utils import secure_filename
import copy

app = Flask(__name__)
CORS(app)

app.config['MAX_CONTENT_LENGTH'] = 30 * 1024 * 1024  # 30 MB upload limit
ALLOWED_EXT = {'.pptx', '.potx'}


# --- Enhanced LLM provider calls ---
def call_llm_split_to_slides(text, guidance, provider, api_key, custom_endpoint=None, max_tokens=1500):
    system_prompt = (
        "You are an expert presentation designer. Given text and optional guidance, create an engaging slide structure. "
        "Analyze the content and create 4-12 slides with compelling titles and well-organized bullet points. "
        "Return ONLY valid JSON in this exact format: "
        "{\"slides\": [{\"title\": \"Slide Title\", \"content\": \"• Point 1\\n• Point 2\\n• Point 3\"}]}"
    )
    
    user_prompt = f"Guidance: {guidance or 'Create a professional presentation'}\n\nContent to structure:\n{text[:20000]}"

    try:
        if provider == 'openai':
            return call_openai(system_prompt, user_prompt, api_key, max_tokens)
        elif provider == 'anthropic':
            return call_anthropic(system_prompt, user_prompt, api_key, max_tokens)
        elif provider == 'gemini':
            return call_gemini(system_prompt, user_prompt, api_key, max_tokens)
        elif provider == 'aipipe':
            return call_aipipe(system_prompt, user_prompt, api_key, max_tokens)
        else:
            raise Exception(f'Unsupported provider: {provider}')
    except Exception as e:
        raise Exception(f'LLM API Error ({provider}): {str(e)}')


def call_openai(system_prompt, user_prompt, api_key, max_tokens):
    url = "https://api.openai.com/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    payload = {
        "model": "gpt-4o-mini",
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        "temperature": 0.3,
        "max_tokens": max_tokens
    }
    
    response = requests.post(url, headers=headers, json=payload, timeout=60)
    response.raise_for_status()
    result = response.json()
    raw_content = result['choices'][0]['message']['content']
    return parse_llm_response(raw_content)


def call_anthropic(system_prompt, user_prompt, api_key, max_tokens):
    url = "https://api.anthropic.com/v1/messages"
    headers = {
        "x-api-key": api_key,
        "Content-Type": "application/json",
        "anthropic-version": "2023-06-01"
    }
    payload = {
        "model": "claude-3-haiku-20240307",
        "max_tokens": max_tokens,
        "system": system_prompt,
        "messages": [
            {
                "role": "user", 
                "content": user_prompt
            }
        ],
        "temperature": 0.3
    }
    
    response = requests.post(url, headers=headers, json=payload, timeout=60)
    response.raise_for_status()
    result = response.json()
    raw_content = result['content'][0]['text']
    return parse_llm_response(raw_content)


def call_gemini(system_prompt, user_prompt, api_key, max_tokens):
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash-latest:generateContent?key={api_key}"
    headers = {"Content-Type": "application/json"}
    payload = {
        "contents": [{
            "parts": [{
                "text": f"{system_prompt}\n\n{user_prompt}"
            }]
        }],
        "generationConfig": {
            "temperature": 0.3,
            "maxOutputTokens": max_tokens
        }
    }
    
    response = requests.post(url, headers=headers, json=payload, timeout=60)
    response.raise_for_status()
    result = response.json()
    raw_content = result['candidates'][0]['content']['parts'][0]['text']
    return parse_llm_response(raw_content)


def call_aipipe(system_prompt, user_prompt, api_key, max_tokens):
    url = "https://aipipe.org/openai/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    payload = {
        "model": "gpt-4o-mini",
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        "temperature": 0.3,
        "max_tokens": max_tokens
    }
    
    response = requests.post(url, headers=headers, json=payload, timeout=60)
    response.raise_for_status()
    result = response.json()
    raw_content = result['choices'][0]['message']['content']
    return parse_llm_response(raw_content)


def parse_llm_response(raw_content):
    """Extract and parse JSON from LLM response"""
    try:
        # Find JSON in the response
        start_idx = raw_content.find('{')
        end_idx = raw_content.rfind('}') + 1
        
        if start_idx == -1 or end_idx == 0:
            raise Exception("No JSON found in response")
        
        json_str = raw_content[start_idx:end_idx]
        parsed = json.loads(json_str)
        
        if 'slides' not in parsed or not isinstance(parsed['slides'], list):
            raise Exception("Invalid JSON structure")
        
        return parsed['slides']
        
    except Exception as e:
        # Fallback: create slides from raw text
        print(f"JSON parsing failed: {e}")
        return create_fallback_slides(raw_content)


def create_fallback_slides(text):
    """Create basic slides when JSON parsing fails"""
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    slides = []
    
    current_title = "Introduction"
    current_content = []
    
    for line in lines:
        if len(line) > 50 and not line.startswith(('•', '-', '*')):
            # Likely a title
            if current_content:
                slides.append({
                    "title": current_title,
                    "content": '\n'.join(current_content)
                })
            current_title = line[:80]  # Limit title length
            current_content = []
        else:
            current_content.append(f"• {line.lstrip('•-* ')}")
    
    # Add final slide
    if current_content:
        slides.append({
            "title": current_title,
            "content": '\n'.join(current_content)
        })
    
    return slides[:10]  # Limit to 10 slides


# --- Enhanced template processing ---
def allowed_file(filename):
    _, ext = os.path.splitext(filename.lower())
    return ext in ALLOWED_EXT


def create_enhanced_presentation(slides_struct, template_prs, template_assets=None):
    """
    FIXED: Create presentation preserving template design
    The key is to work directly with the template presentation instead of creating new one
    """
    
    # Start by clearing existing slides from template (keep first slide as reference)
    # We'll work backwards to avoid index shifting
    slide_count = len(template_prs.slides)
    
    # Remove all slides except the first one (we'll use it as reference)
    for i in range(slide_count - 1, 0, -1):
        try:
            rId = template_prs.slides._sldIdLst[i].rId
            template_prs.part.drop_rel(rId)
            del template_prs.slides._sldIdLst[i]
        except Exception as e:
            print(f"Warning: Could not remove slide {i}: {e}")
    
    # If we still have slides, clear the content of the first one, otherwise add a new one
    if len(template_prs.slides) > 0:
        # Clear the first slide but keep its layout
        first_slide = template_prs.slides[0]
        slide_layout = first_slide.slide_layout
        
        # Remove the first slide
        try:
            rId = template_prs.slides._sldIdLst[0].rId
            template_prs.part.drop_rel(rId)
            del template_prs.slides._sldIdLst[0]
        except:
            pass
    
    # Now add new slides with content using available layouts
    available_layouts = template_prs.slide_layouts
    
    for i, slide_data in enumerate(slides_struct):
        try:
            # Cycle through available layouts, prefer content layouts (usually index 1-6)
            if len(available_layouts) > 1:
                # Skip title slide layout (usually index 0) for content slides after the first
                layout_idx = 1 + (i % max(1, len(available_layouts) - 1))
                if layout_idx >= len(available_layouts):
                    layout_idx = 1 if len(available_layouts) > 1 else 0
            else:
                layout_idx = 0
            
            slide_layout = available_layouts[layout_idx]
            slide = template_prs.slides.add_slide(slide_layout)
            
            # Add title
            title = slide_data.get('title', f'Slide {i+1}')
            
            # Try to find title placeholder
            title_added = False
            for shape in slide.placeholders:
                try:
                    if hasattr(shape, 'text') and shape.placeholder_format.idx == 0:  # Title placeholder
                        shape.text = title
                        title_added = True
                        break
                except:
                    continue
            
            # If no title placeholder found, try shapes.title
            if not title_added and hasattr(slide.shapes, 'title') and slide.shapes.title:
                try:
                    slide.shapes.title.text = title
                    title_added = True
                except:
                    pass
            
            # Add content
            content = slide_data.get("content", "").strip()
            if content:
                content_added = False
                
                # Try to find content placeholder (usually idx 1 or 10)
                for shape in slide.placeholders:
                    try:
                        if (hasattr(shape, 'text_frame') and 
                            shape.placeholder_format.idx != 0 and  # Not title
                            shape.text_frame is not None):
                            
                            text_frame = shape.text_frame
                            text_frame.clear()  # Clear existing content
                            
                            # Add content lines
                            lines = [line.strip() for line in content.split('\n') if line.strip()]
                            
                            for j, line in enumerate(lines):
                                # Clean bullet formatting
                                clean_line = line.lstrip('•-* ').strip()
                                if not clean_line:
                                    continue
                                
                                if j == 0:
                                    p = text_frame.paragraphs[0]
                                else:
                                    p = text_frame.add_paragraph()
                                
                                p.text = clean_line
                                
                                # Set appropriate bullet level
                                if line.startswith(('•', '-', '*')):
                                    p.level = 1
                                else:
                                    p.level = 0
                            
                            content_added = True
                            break
                    except Exception as e:
                        print(f"Error adding content to placeholder: {e}")
                        continue
                
                # Fallback: add textbox if no content placeholder found
                if not content_added:
                    try:
                        # Calculate position to avoid title area
                        top_pos = Inches(2.5) if title_added else Inches(1.5)
                        textbox = slide.shapes.add_textbox(
                            Inches(0.5), top_pos, Inches(9), Inches(5)
                        )
                        text_frame = textbox.text_frame
                        text_frame.word_wrap = True
                        
                        lines = [line.strip() for line in content.split('\n') if line.strip()]
                        for j, line in enumerate(lines):
                            clean_line = line.lstrip('•-* ').strip()
                            if not clean_line:
                                continue
                            
                            if j == 0:
                                p = text_frame.paragraphs[0]
                            else:
                                p = text_frame.add_paragraph()
                            
                            p.text = clean_line
                            if line.startswith(('•', '-', '*')):
                                p.level = 1
                    except Exception as e:
                        print(f"Error creating fallback textbox: {e}")
        
        except Exception as e:
            print(f"Error creating slide {i+1}: {e}")
            continue
    
    # Return the modified template presentation
    output_stream = io.BytesIO()
    template_prs.save(output_stream)
    output_stream.seek(0)
    return output_stream


# --- API Endpoints ---
@app.route('/generate', methods=['POST'])
def generate():
    try:
        # Validate inputs
        text = request.form.get('text', '').strip()
        provider = request.form.get('provider', 'openai')
        api_key = request.form.get('api_key', '').strip()
        guidance = request.form.get('guidance', '').strip()

        if not text:
            return jsonify({'error': 'Text content is required'}), 400
        
        if not api_key:
            return jsonify({'error': 'API key is required'}), 400

        if 'template' not in request.files:
            return jsonify({'error': 'Template file is required'}), 400

        template_file = request.files['template']
        if not template_file or not allowed_file(template_file.filename):
            return jsonify({'error': 'Invalid template file. Please upload .pptx or .potx'}), 400

        # Process template
        template_stream = io.BytesIO(template_file.read())
        template_prs = Presentation(template_stream)

        # Generate slide structure using LLM
        slides_structure = call_llm_split_to_slides(text, guidance, provider, api_key)
        
        if not slides_structure:
            return jsonify({'error': 'Failed to generate slide structure'}), 500

        # Create presentation with proper template preservation
        output_stream = create_enhanced_presentation(slides_structure, template_prs)

        return send_file(
            output_stream,
            as_attachment=True,
            download_name='generated_presentation.pptx',
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except Exception as e:
        print(f"Error in generate endpoint: {str(e)}")
        print(traceback.format_exc())
        return jsonify({'error': str(e)}), 500


@app.route('/test-api', methods=['POST'])
def test_api():
    """Test endpoint to validate API keys"""
    try:
        provider = request.json.get('provider')
        api_key = request.json.get('api_key')
        
        if not provider or not api_key:
            return jsonify({'valid': False, 'error': 'Missing provider or API key'}), 400
        
        # Test with a simple prompt
        test_slides = call_llm_split_to_slides(
            "Test content for API validation", 
            "Create a simple test slide", 
            provider, 
            api_key,
            max_tokens=100
        )
        
        return jsonify({'valid': True, 'message': 'API key is valid'})
        
    except Exception as e:
        return jsonify({'valid': False, 'error': str(e)}), 400


@app.route('/health')
def health():
    return jsonify({'status': 'healthy'})


@app.route('/')
def index():
    try:
        with open("index.html", encoding="utf-8") as f:
            html_content = f.read()
        return render_template_string(html_content)
    except FileNotFoundError:
        return "Template file not found", 404


@app.route('/styles.css')
def serve_css():
    try:
        with open('styles.css', 'r', encoding='utf-8') as f:
            css_content = f.read()
        return css_content, 200, {'Content-Type': 'text/css; charset=utf-8'}
    except FileNotFoundError:
        return "CSS file not found", 404



if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    host = os.environ.get('HOST', '0.0.0.0')
    debug = os.environ.get('FLASK_ENV') != 'production'
    
    app.run(debug=debug, host=host, port=port)